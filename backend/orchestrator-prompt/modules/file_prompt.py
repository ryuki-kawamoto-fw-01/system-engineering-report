import base64
import csv
import io
import logging

import chardet
import docx
import fitz
import openpyxl
import pptx
from msg_parser import MsOxMessage
from openai import BadRequestError

from .errors.custom_errors import LLMSpecificError


def _try_decode(file_stream: bytes, detected_encoding: str | None) -> str:
    """複数のエンコーディングを試行してデコードする"""
    if detected_encoding and detected_encoding.lower() in ["shift_jis", "sjis"]:
        encodings_to_try = ["shift_jis", "utf-8", "cp932", "latin-1"]
    elif detected_encoding:
        encodings_to_try = [detected_encoding, "utf-8", "shift_jis", "cp932", "latin-1"]
    else:
        encodings_to_try = ["utf-8", "shift_jis", "cp932", "latin-1"]
    
    for enc in encodings_to_try:
        try:
            content = file_stream.decode(enc)
            logging.info(f"Successfully decoded with: {enc}")
            return content
        except (UnicodeDecodeError, LookupError):
            continue
    
    raise ValueError("ファイルのエンコーディングを判定できませんでした")


def get_file_content(file_stream: bytes, file_extension: str):
    try:
        if file_extension == "txt":
            result = chardet.detect(file_stream)
            encoding = result["encoding"]
            logging.info(f"Detected encoding: {encoding}")
            file_content = _try_decode(file_stream, encoding)
        elif file_extension == "pdf":
            pdf_document = fitz.open(stream=file_stream, filetype="pdf")
            extracted_text = ""
            for page in pdf_document:
                extracted_text += page.get_text()
            file_content = extracted_text
        elif file_extension == "docx":
            doc = docx.Document(io.BytesIO(file_stream))
            extracted_text = ""
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
            file_content = extracted_text
        elif file_extension == "xlsx":
            workbook = openpyxl.load_workbook(io.BytesIO(file_stream))
            extracted_text = ""
            for sheet in workbook:
                for row in sheet.iter_rows(values_only=True):
                    extracted_text += (
                        "\t".join([str(cell) for cell in row if cell is not None])
                        + "\n"
                    )
            file_content = extracted_text
        elif file_extension == "pptx":
            presentation = pptx.Presentation(io.BytesIO(file_stream))
            extracted_text = ""
            for slide in presentation.slides:
                # 図形を位置に基づいてソート(スライド内の図形を上から↓、左から右の順序でソート)
                # ユーザーが視覚的に認識する順序に近い形でテキストを取得
                shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
                for shape in shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
            file_content = extracted_text
        elif file_extension == "csv":
            result = chardet.detect(file_stream)
            encoding = result["encoding"]
            logging.info(f"Detected encoding for CSV: {encoding}")
            file_text = _try_decode(file_stream, encoding)
            csv_reader = csv.reader(io.StringIO(file_text))
            extracted_text = ""
            for row in csv_reader:
                extracted_text += ",".join(row) + "\n"
            file_content = extracted_text
        elif file_extension == "msg":
            msg_obj = MsOxMessage(io.BytesIO(file_stream))
            msg_properties_dict = msg_obj.get_properties()
            extracted_text = (
                "件名: "
                + msg_properties_dict.get("Subject", "")
                + "\n\n"
                + msg_properties_dict.get("Body", "")
            )
            file_content = extracted_text
        else:
            file_content = ""
            logging.info(f"Unsupported file extension: {file_extension}")

    except Exception as e:
        file_content = ""
        logging.error(f"Error decoding file content: {e}")
        if file_extension == "msg":
            raise LLMSpecificError(
                "E_B3_00120",
                value="メッセージ"
            )
        else:
            raise LLMSpecificError(
                "E_B3_00130",
                value="メッセージ"
            )

    return file_content

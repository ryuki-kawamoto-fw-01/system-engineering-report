import io
import logging
from typing import List

from pptx import Presentation


# パワポファイル(バイト列)をスライドごとに分割
def split_pptx(input_file: bytes) -> List[bytes]:
    outputs: List[bytes] = []

    try:
        # プレゼンテーションを読み込む
        prs = Presentation(io.BytesIO(input_file))
    except Exception:
        logging.error("ファイルを読み込めません。暗号化されている可能性があります")
        raise

    total_slides = len(prs.slides)

    for i in range(total_slides):
        # 元のプレゼンテーションをコピー
        new_prs = Presentation(io.BytesIO(input_file))

        # 他のスライドを削除
        slide_indices_to_remove = [j for j in range(total_slides) if j != i]
        for index in sorted(slide_indices_to_remove, reverse=True):
            rId = new_prs.slides._sldIdLst[index].rId
            new_prs.part.drop_rel(rId)
            del new_prs.slides._sldIdLst[index]

        # 新しいプレゼンテーションをメモリ内に保存
        new_pptx_io = io.BytesIO()
        new_prs.save(new_pptx_io)
        new_pptx_io.seek(0)
        outputs.append(new_pptx_io.getvalue())

    return outputs

import base64
import csv
import io
import logging

import chardet
import docx
import fitz
import openpyxl
import pptx
from msg_parser import MsOxMessage


def _try_decode(content: bytes, detected_encoding: str | None) -> str:
    """複数のエンコーディングを試行してデコードする"""
    if detected_encoding and detected_encoding.lower() in ["shift_jis", "sjis"]:
        encodings_to_try = ["shift_jis", "utf-8", "cp932", "latin-1"]
    elif detected_encoding:
        encodings_to_try = [detected_encoding, "utf-8", "shift_jis", "cp932", "latin-1"]
    else:
        encodings_to_try = ["utf-8", "shift_jis", "cp932", "latin-1"]

    for enc in encodings_to_try:
        try:
            result = content.decode(enc)
            logging.info(f"Successfully decoded with: {enc}")
            return result
        except (UnicodeDecodeError, LookupError):
            continue

    raise ValueError("ファイルのエンコーディングを判定できませんでした")


def get_file_content(file_content: str, file_extension: str) -> str:
    """
    ファイルの拡張子に基づいてファイルの内容をデコードし、テキスト内容を返します。
    引数：
        file_content (str): Base64でエンコードされたファイル内容。
        file_extension (str): ファイルの拡張子（例：'txt', 'pdf', 'docx', 'xlsx', 'pptx', 'csv'）。
    戻り値：
        str: ファイルのデコードされたテキスト内容。
    """

    try:
        decoded_content = base64.b64decode(file_content)
        if file_extension == "txt":
            result = chardet.detect(decoded_content)
            encoding = result["encoding"]
            logging.info(f"Detected encoding: {encoding}")
            file_content = _try_decode(decoded_content, encoding)
        elif file_extension == "pdf":
            pdf_document = fitz.open(stream=decoded_content, filetype="pdf")
            extracted_text = ""
            for page in pdf_document:
                extracted_text += page.get_text()
            file_content = extracted_text
        elif file_extension == "docx":
            doc = docx.Document(io.BytesIO(decoded_content))
            extracted_text = ""
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
            file_content = extracted_text
        elif file_extension == "xlsx":
            workbook = openpyxl.load_workbook(io.BytesIO(decoded_content))
            extracted_text = ""
            for sheet in workbook:
                for row in sheet.iter_rows(values_only=True):
                    extracted_text += (
                        "\t".join([str(cell) for cell in row if cell is not None])
                        + "\n"
                    )
            file_content = extracted_text
        elif file_extension == "pptx":
            presentation = pptx.Presentation(io.BytesIO(decoded_content))
            extracted_text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
            file_content = extracted_text
        elif file_extension == "csv":
            result = chardet.detect(decoded_content)
            encoding = result["encoding"]
            logging.info(f"Detected encoding for CSV: {encoding}")
            decoded_content = _try_decode(decoded_content, encoding)
            csv_reader = csv.reader(io.StringIO(decoded_content))
            extracted_text = ""
            for row in csv_reader:
                extracted_text += ",".join(row) + "\n"
            file_content = extracted_text
        elif file_extension == "msg":
            msg_obj = MsOxMessage(io.BytesIO(decoded_content))
            msg_properties_dict = msg_obj.get_properties()
            extracted_text = (
                "件名: "
                + msg_properties_dict.get("Subject", "")
                + "\n\n"
                + msg_properties_dict.get("Body", "")
            )
            file_content = extracted_text
        else:
            file_content = ""
            logging.info(f"Unsupported file extension: {file_extension}")

    except Exception as e:
        file_content = ""
        logging.error(f"Error decoding file content: {e}")

    return file_content
